"""
Proposal text extraction — no AI, pure document parsing.
Returns an ordered list of paragraph strings (blank lines preserved as "").
Claude Code does all the reasoning; this module only pulls raw text.
"""
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_proposal_text(filepath: str) -> list[str]:
    """
    Extract text from a .docx, .pdf, or .pptx proposal.
    Returns ordered list of paragraph strings.
    """
    path = Path(filepath)
    ext = path.suffix.lower()

    if ext == ".docx":
        return _docx(path)
    elif ext == ".pdf":
        return _pdf(path)
    elif ext == ".pptx":
        return _pptx(path)
    else:
        raise ValueError(f"Unsupported proposal format '{ext}'. Use .docx, .pdf, or .pptx.")


def _docx(path: Path) -> list[str]:
    from docx import Document
    doc = Document(path)
    paras = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paras.append(" | ".join(cells))
    return paras


def _pdf(path: Path) -> list[str]:
    import pdfplumber
    paras: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                paras.extend(text.splitlines())
            else:
                logger.warning("Page %d has no extractable text — trying OCR.", page.page_number)
                try:
                    import pytesseract
                    img = page.to_image(resolution=150).original
                    paras.extend(pytesseract.image_to_string(img).splitlines())
                except Exception as e:
                    logger.warning("OCR failed for page %d: %s", page.page_number, e)

    stripped = [p.strip() for p in paras]
    if not any(stripped):
        raise ValueError(
            "No text extracted from the PDF. "
            "If it is a scanned document, install Tesseract for OCR support."
        )
    return stripped


def _pptx(path: Path) -> list[str]:
    from pptx import Presentation
    paras: list[str] = []
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    paras.append(p.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        paras.append(" | ".join(cells))
    return paras
