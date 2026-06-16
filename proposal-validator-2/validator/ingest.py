"""
Source document ingestion and indexing.
Parses .xlsx/.xls, .csv, .docx/.doc, .pdf, .pptx/.ppt files into a flat list of
searchable text chunks, each tagged with provenance metadata.

Each file is parsed in a worker thread with a 60-second timeout so a single
bad file (e.g. a corrupt PDF) can never hang the whole run.
"""
import csv
import logging
from collections.abc import Callable
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeout
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".xlsx", ".xls", ".csv", ".docx", ".doc", ".pdf", ".pptx", ".ppt"}
FILE_TIMEOUT = 60  # seconds per file before we skip it


def ingest_sources(
    folder_path: str,
    on_file: Callable[[str, int, int], None] | None = None,
) -> list[dict]:
    """
    Walk a folder, parse every supported file, return a flat list of chunks.

    on_file(filename, file_num, total_files) is called after each file attempt
    so the caller can stream progress to the UI.
    """
    folder = Path(folder_path)
    if not folder.exists() or not folder.is_dir():
        raise ValueError(f"Source folder does not exist: {folder_path}")

    files = [
        f for f in sorted(folder.iterdir())
        if not f.is_dir() and f.suffix.lower() in SUPPORTED_EXTENSIONS
    ]
    total = len(files)
    chunks = []

    for idx, file_path in enumerate(files, 1):
        ext = file_path.suffix.lower()
        parser = {
            ".xlsx": _parse_excel, ".xls": _parse_excel,
            ".csv":  _parse_csv,
            ".docx": _parse_word,  ".doc":  _parse_word,
            ".pdf":  _parse_pdf,
            ".pptx": _parse_pptx, ".ppt":  _parse_pptx,
        }.get(ext)

        if parser:
            new_chunks = _parse_with_timeout(parser, file_path)
            chunks.extend(new_chunks)
            logger.info("[%d/%d] %s -> %d chunks", idx, total, file_path.name, len(new_chunks))
        else:
            logger.info("Skipping unsupported file: %s", file_path.name)

        if on_file:
            on_file(file_path.name, idx, total)

    logger.info("Ingested %d chunks total from %s", len(chunks), folder_path)
    return chunks


def _parse_with_timeout(parser, file_path: Path) -> list[dict]:
    """Run parser in a thread; return [] if it times out or raises."""
    with ThreadPoolExecutor(max_workers=1) as ex:
        future = ex.submit(parser, file_path)
        try:
            return future.result(timeout=FILE_TIMEOUT)
        except FutureTimeout:
            logger.warning("Timeout (>%ds) parsing %s -- skipping.", FILE_TIMEOUT, file_path.name)
            return []
        except Exception as exc:
            logger.warning("Failed to parse %s: %s -- skipping.", file_path.name, exc)
            return []


# ---------------------------------------------------------------------------
# Per-format parsers
# ---------------------------------------------------------------------------

def _chunk(filename, filetype, location, sheet, page_or_slide, row, text):
    return {
        "filename": filename,
        "filetype": filetype,
        "location": location,
        "sheet": sheet,
        "page_or_slide": page_or_slide,
        "row": row,
        "text": text.strip(),
    }


def _parse_excel(file_path: Path) -> list[dict]:
    chunks = []
    ext = file_path.suffix.lower()
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            headers = [
                str(c.value) if c.value is not None else f"Col{i+1}"
                for i, c in enumerate(next(ws.iter_rows(min_row=1, max_row=1)))
            ]
            for row_idx, row in enumerate(ws.iter_rows(min_row=1, values_only=True), start=1):
                parts = []
                for col_idx, val in enumerate(row):
                    if val is not None:
                        hdr = headers[col_idx] if col_idx < len(headers) else f"Col{col_idx+1}"
                        parts.append(f"{hdr}: {val}")
                if parts:
                    loc = f"Sheet: {sheet_name}, Row: {row_idx}"
                    chunks.append(_chunk(file_path.name, "xlsx", loc, sheet_name, None, row_idx, " | ".join(parts)))
    except Exception:
        try:
            import pandas as pd
            engine = "xlrd" if ext == ".xls" else "openpyxl"
            xl = pd.ExcelFile(file_path, engine=engine)
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name, dtype=str)
                for row_idx, (_, row) in enumerate(df.iterrows(), start=2):
                    parts = [f"{col}: {val}" for col, val in row.items() if str(val) not in ("nan", "")]
                    if parts:
                        loc = f"Sheet: {sheet_name}, Row: {row_idx}"
                        chunks.append(_chunk(file_path.name, "xlsx", loc, sheet_name, None, row_idx, " | ".join(parts)))
        except Exception as exc2:
            logger.warning("pandas fallback failed for %s: %s", file_path.name, exc2)
    return chunks


def _parse_csv(file_path: Path) -> list[dict]:
    chunks = []
    for encoding in ("utf-8-sig", "latin-1"):
        try:
            with open(file_path, newline="", encoding=encoding) as f:
                reader = csv.DictReader(f)
                for row_idx, row in enumerate(reader, start=2):
                    parts = [f"{k}: {v}" for k, v in row.items() if v and v.strip()]
                    if parts:
                        chunks.append(_chunk(file_path.name, "csv", f"Row: {row_idx}", None, None, row_idx, " | ".join(parts)))
            break
        except UnicodeDecodeError:
            continue
    return chunks


def _parse_word(file_path: Path) -> list[dict]:
    chunks = []
    ext = file_path.suffix.lower()
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            for para_idx, para in enumerate(doc.paragraphs, start=1):
                text = para.text.strip()
                if text:
                    chunks.append(_chunk(file_path.name, "docx", f"Paragraph: {para_idx}", None, None, para_idx, text))
            for tbl_idx, table in enumerate(doc.tables, start=1):
                for row_idx, row in enumerate(table.rows, start=1):
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        chunks.append(_chunk(file_path.name, "docx", f"Table {tbl_idx}, Row {row_idx}", None, None, row_idx, " | ".join(cells)))
            return chunks
        except Exception as exc:
            logger.warning("python-docx failed for %s: %s", file_path.name, exc)
    try:
        import docx2txt
        text = docx2txt.process(str(file_path))
        for line_idx, line in enumerate(text.splitlines(), start=1):
            line = line.strip()
            if line:
                chunks.append(_chunk(file_path.name, "doc", f"Line: {line_idx}", None, None, line_idx, line))
    except Exception as exc:
        logger.warning("docx2txt failed for %s: %s", file_path.name, exc)
    return chunks


def _parse_pdf(file_path: Path) -> list[dict]:
    chunks = []
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                pg_num = page.page_number
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    for line in page_text.splitlines():
                        line = line.strip()
                        if line:
                            chunks.append(_chunk(file_path.name, "pdf", f"Page: {pg_num}", None, pg_num, None, line))
                else:
                    try:
                        import pytesseract
                        img = page.to_image(resolution=150).original
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            chunks.append(_chunk(file_path.name, "pdf", f"Page: {pg_num} (OCR)", None, pg_num, None, ocr_text.strip()))
                    except Exception as ocr_exc:
                        logger.warning("OCR failed for %s page %d: %s", file_path.name, pg_num, ocr_exc)
                for tbl_idx, table in enumerate(page.extract_tables() or [], start=1):
                    for row_idx, row in enumerate(table, start=1):
                        cells = [str(c).strip() for c in row if c and str(c).strip()]
                        if cells:
                            chunks.append(_chunk(file_path.name, "pdf", f"Page: {pg_num}, Table {tbl_idx}, Row {row_idx}", None, pg_num, row_idx, " | ".join(cells)))
    except Exception as exc:
        logger.warning("pdfplumber failed for %s: %s", file_path.name, exc)
    return chunks


def _parse_pptx(file_path: Path) -> list[dict]:
    chunks = []
    if file_path.suffix.lower() == ".ppt":
        logger.warning("Legacy .ppt not supported: %s -- skipping.", file_path.name)
        return chunks
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            chunks.append(_chunk(file_path.name, "pptx", f"Slide: {slide_idx}, Shape: {shape.name}", None, slide_idx, None, text))
                if shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows, start=1):
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            chunks.append(_chunk(file_path.name, "pptx", f"Slide: {slide_idx}, Table Row: {row_idx}", None, slide_idx, row_idx, " | ".join(cells)))
    except Exception as exc:
        logger.warning("python-pptx failed for %s: %s", file_path.name, exc)
    return chunks
